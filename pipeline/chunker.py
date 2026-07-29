import hashlib
import re
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import List, Optional

from config.settings import CHUNK_MAX_WORDS, CHUNK_OVERLAP_WORDS

@dataclass
class Chunk:
    id: int
    file: str
    document_id: str
    doc_title: str
    section: str
    document_type: str
    page_number: int
    start_line: int
    end_line: int
    upload_timestamp: str
    chunk_id: str
    text: str
    department: Optional[str] = None
    system: Optional[str] = None
    service: Optional[str] = None
    environment: Optional[str] = None
    operating_system: Optional[str] = None
    tags: List[str] = field(default_factory=list)

def slugify(text: str) -> str:
    text = text.lower()
    text = re.sub(r'[^a-z0-9]+', '-', text)
    return text.strip('-')

def is_atomic_line(line: str, in_code_block: bool) -> bool:
    if in_code_block:
        return True
    line_s = line.strip()
    if line_s.startswith('```') or line_s.startswith('~~~'):
        return True
    prefixes = ('$', '>', '#!', 'kubectl', 'docker', 'helm', 'aws', 'git', 'python', 'pip', 'apt', 'yum', 'rpm', 'systemctl', 'service', 'chmod', 'chown', 'mkdir', 'cp', 'mv', 'rm', 'cat', 'grep', 'sed', 'awk', 'curl', 'wget', 'ps', 'kill', 'export', 'source', '.')
    if any(line_s.startswith(p) for p in prefixes):
        return True
    if '|' in line_s:
        return True
    if re.match(r'^\d+\.', line_s):
        return True
    admonitions = ('WARNING', 'CAUTION', 'NOTE', 'IMPORTANT', 'DANGER', '**WARNING**', '**CAUTION**')
    if any(line_s.startswith(a) for a in admonitions):
        return True
    return False

def chunk_file(path: Path, start_id: int) -> List[Chunk]:
    chunks = []
    try:
        content = path.read_bytes()
        doc_id = hashlib.sha256(content).hexdigest()
        
        doc_type = path.suffix.strip('.').lower()
        if doc_type == 'pdf':
            from pdfminer.high_level import extract_text as pdf_extract_text
            text = pdf_extract_text(str(path))
            doc_type = 'pdf'
        elif doc_type == 'docx':
            import docx
            doc = docx.Document(str(path))
            text_runs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_runs.append(para.text.strip())
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        text_runs.append(" | ".join(row_text))
            text = '\n'.join(text_runs)
            doc_type = 'docx'
        elif doc_type == 'pptx':
            from pptx import Presentation
            prs = Presentation(str(path))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                text_runs.append(" | ".join(row_text))
            text = '\n'.join(text_runs)
            doc_type = 'pptx'
        else:
            text = content.decode('utf-8', errors='ignore')
            if doc_type == 'md': doc_type = 'markdown'
            else: doc_type = 'unknown'
            
    except Exception as e:
        print(f"[chunker] Failed to parse {path.name}: {e}")
        return chunks

    lines = text.split('\n')

    doc_title = path.stem
    section = "General"
    department, system, service, env, os_val = None, None, None, None, None
    tags = []

    # Simple frontmatter parsing
    if lines and lines[0].strip() == '---':
        for i in range(1, len(lines)):
            if lines[i].strip() == '---':
                break
            if ':' in lines[i]:
                k, v = lines[i].split(':', 1)
                k = k.strip().lower()
                v = v.strip()
                if k == 'department': department = v
                elif k == 'system': system = v
                elif k == 'service': service = v
                elif k == 'environment': env = v
                elif k == 'operating_system': os_val = v
                elif k == 'tags': tags = [t.strip() for t in v.split(',')]
    
    current_chunk_text = []
    start_line = 1
    total_lines = len(lines) if len(lines) > 0 else 1
    chunk_idx = 0
    in_code = False

    # Simple sequential chunker adhering to constraints
    for i, line in enumerate(lines):
        line_num = i + 1
        if line.strip().startswith('# ') and not in_code:
            doc_title = line.strip()[2:].strip()
            continue
        if line.strip().startswith('## ') and not in_code:
            section = line.strip()[3:].strip()
            continue
        
        if line.strip().startswith('```') or line.strip().startswith('~~~'):
            in_code = not in_code
        
        current_chunk_text.append(line)
        word_count = len(' '.join(current_chunk_text).split())
        
        if word_count >= CHUNK_MAX_WORDS and not in_code:
            # Create chunk
            page_number = int(start_line / total_lines * 10) + 1
            chunks.append(Chunk(
                id=start_id, file=path.name, document_id=doc_id, doc_title=doc_title, section=section,
                document_type=doc_type, page_number=page_number, start_line=start_line, end_line=line_num,
                upload_timestamp=datetime.utcnow().isoformat() + "Z",
                chunk_id=f"{doc_id[:8]}-{slugify(section)}-{chunk_idx}",
                text='\n'.join(current_chunk_text),
                department=department, system=system, service=service, environment=env, operating_system=os_val, tags=tags
            ))
            start_id += 1
            chunk_idx += 1
            
            # Implement overlap by keeping the last few lines that approximate CHUNK_OVERLAP_WORDS
            overlap_lines = []
            overlap_words = 0
            for prev_line in reversed(current_chunk_text):
                if overlap_words >= CHUNK_OVERLAP_WORDS:
                    break
                overlap_lines.insert(0, prev_line)
                overlap_words += len(prev_line.split())
                
            current_chunk_text = overlap_lines
            start_line = line_num - len(overlap_lines) + 1

    if current_chunk_text:
        page_number = int(start_line / total_lines * 10) + 1
        chunks.append(Chunk(
            id=start_id, file=path.name, document_id=doc_id, doc_title=doc_title, section=section,
            document_type=doc_type, page_number=page_number, start_line=start_line, end_line=total_lines,
            upload_timestamp=datetime.utcnow().isoformat() + "Z",
            chunk_id=f"{doc_id[:8]}-{slugify(section)}-{chunk_idx}",
            text='\n'.join(current_chunk_text),
            department=department, system=system, service=service, environment=env, operating_system=os_val, tags=tags
        ))
    return chunks
